# Bibliotecas usadas
import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from groq import Groq
import time

# Definição da função que extrai o texto em pdf
def extrair_texto_pdf(caminho_arq):
    # Carrega o ficheiro pdf na memória
    leitor = PdfReader(caminho_arq)
    texto = ""
    # Percorre cada página do PDF e vai acumulando texto extraído 
    for pagina in leitor.pages:
        texto += pagina.extract_text() + "\n"
    return texto

# Definição da função que extrai o texto em word
def extrair_texto_word(caminho_arq):
    # Carrega o ficheiro word na memória 
    doc = Document(caminho_arq)
    # Usa uma list comprhension para extrair o texto de cada parágrafo do documento 
    texto = [paragrafo.text for paragrafo in doc.paragraphs]
    # Junta todos os parágrafos recolhidos e quebra linhas entre eles 
    return "\n".join(texto)

# Definição da função que extrai o texto em powerpoint 
def extrair_texto_powerpoint(caminho_arq):
    # Carrega o ficheiro PowerPoint na memória 
    prs = Presentation(caminho_arq)
    texto = []
    # Percorre cada slide da apresentação 
    for slide in prs.slides:
        # Percorre cada elemento dentro do slide atual
        for forma in slide.shapes:
            # Verifica se o elemento possui a propriedade de texto 
            if hasattr(forma, "text"):
                texto.append(forma.text)
    # Junta todos os textos recolhidos dentro de uma única string 
    return "\n".join(texto)

# Definição da função que identifica a extensão, lê e gera o resumo  
def ler_resumir(caminho_arq):
    # Separa o caminho da extensão para direcionar o ficheiro à função da extração correta 
    _, extensao = os.path.splitext(caminho_arq.lower())
    texto_completo = ""

    print(f"\nLendo o arquivo: {caminho_arq}...")
    # Direciona o ficheiro à função de extração correta
    if extensao == ".pdf":
        texto_completo = extrair_texto_pdf(caminho_arq)
    elif extensao == ".docx":
        texto_completo = extrair_texto_word(caminho_arq)
    elif extensao == ".pptx":
        texto_completo = extrair_texto_powerpoint(caminho_arq)
    # Caso o formato não seja nenhum dos três acima, interrompe a execução da função
    else:
        print("Formato de arquivo não suportado...")
        return
    # Verifica se o texto extraído está vazio
    if not texto_completo.strip():
        print("Não foi possivel extrair nenhum texto deste arquivo.")
        return
    
    print("Gerando o resumo com IA...")
    # Inicializa o cliente da Groq 
    client = Groq()

    # Cria o prompt para o modelo de IA, limitando o texto aos primeiros 15.000 caracteres
    prompt = f"""
    Atue como um assistente de pesquisa. Faça um resumo estruturado e claro do texto abaixo, 
    destacando os pontos principais, objetivos e conclusões em tópicos (bullet points).
    
    Texto:
    {texto_completo[:15000]}  
    """
    # Envia o prompt para o modelo IA configurado na Groq
    resposta = client.chat.completions.create(
        model='llama-3.3-70b-versatile',
        messages=[
            {"role": "user", "content": prompt}
        ]
    )

    # Exibe visualmente o cabeçalho do resumo no terminal
    print("\n" + "="*40)
    print("==== RESUMO DO DOCUMENTO ====")
    print("="*40)
    # Imprime apenas o texto final e retorna pela IA 
    print(resposta.choices[0].message.content)
# Garante que o bloco abaixo só execute se o script for executado diretamente
if __name__ == "__main__":
    print("\nGuarde o arquivo numa pasta para eu acessar !\n")
    # Pede o caminho da pasta com o nome do ficheiro ao utiliador
    caminho_documento = input("Qual a pasta onde se encontra esse arquivo ? ").replace('\\', '/')
    nome_arquivo = input("Qual o arquivo que gostava que eu resumisse ?").strip()
    # Junta o caminho da pasta com o nome do ficheiro de forma segura para criar o caminho completo 
    caminho_completo = os.path.join(caminho_documento, nome_arquivo)

    print("\nVerificando o arquivo...")
    # Delay interativo de 1 segundo
    time.sleep(1)
    # Verifica se o ficheiro realmente existe no caminho antes de tentar abri-lo
    if os.path.exists(caminho_completo):
        ler_resumir(caminho_completo)
    # Se não existir avisa ao utilizador que o ficheiro não foi encontrado
    else:
        print(f"Erro: O arquivo '{nome_arquivo}' não foi encontrado...")